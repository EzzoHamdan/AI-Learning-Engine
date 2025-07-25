import streamlit as st
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import json
import re
from dotenv import load_dotenv

# Auto-configuration for easy setup
try:
    from auto_config import ensure_env_file_exists, get_setup_status
    ensure_env_file_exists()  # Create .env if missing
    ready, status_message, available_providers = get_setup_status()
except ImportError:
    # Fallback if auto_config is not available
    ready, status_message, available_providers = True, "Manual configuration", []

# Load environment variables from .env file
load_dotenv()

# Import custom modules
from config import AppConfig, QuizConfig, OpenAIConfig, GoogleAIConfig, LocalAIConfig, DIFFICULTY_CONFIG, SCORING_CONFIG
from logger import setup_logging
from session_manager import SessionManager
from ai_client_factory import AIClientFactory
from local_ai_client import is_ollama_running, list_available_models
from study_materials_generator import StudyMaterialsGenerator

# Setup logging
logger = setup_logging()

# Initialize configuration
app_config = AppConfig()
quiz_config = QuizConfig()
openai_config = OpenAIConfig()
google_ai_config = GoogleAIConfig()
local_ai_config = LocalAIConfig()

# Initialize session manager and AI client factory
session_manager = SessionManager()
ai_factory = AIClientFactory(session_manager)

# Get AI client with graceful error handling
client, ai_provider, client_successful = ai_factory.get_working_client()

# Debug: Print configuration values
if app_config.DEBUG_MODE:
    st.write("**Debug - Configuration Status:**")
    st.write(f"Selected Provider: {st.session_state.ai_provider}")
    st.write(f"Active Provider: {ai_provider}")
    st.write(f"Client Status: {'✅ Working' if client_successful else '❌ Error Mode'}")
    st.write(f"Provider Status: {st.session_state.provider_status}")

# Show provider status
if client_successful:
    if app_config.DEBUG_MODE:
        st.success(f"✅ Successfully initialized: {ai_provider}")
else:
    st.warning(f"⚠️ AI Provider Issue: {ai_provider}")
    st.info("💡 The app will continue to work. Please configure a working AI provider in the sidebar to generate quizzes.")

def get_model_info(model_name):
    return model_name

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join(para.text for para in doc.paragraphs)

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def summarize_text(text):
    # Check if we have a working client
    if not client_successful:
        st.error("❌ No working AI provider available for text summarization.")
        return text  # Return original text if no AI available
        
    # Use appropriate model based on provider
    if st.session_state.ai_provider == "Local AI (Ollama)":
        # Use selected model from session state if available, otherwise fall back to config
        model = getattr(st.session_state, 'selected_local_model', local_ai_config.MODEL_NAME)
        temperature = local_ai_config.SUMMARY_TEMPERATURE
    elif st.session_state.ai_provider == "Google AI":
        model = google_ai_config.CHAT_MODEL
        temperature = google_ai_config.SUMMARY_TEMPERATURE
    else:
        model = "gpt-3.5-turbo"
        temperature = 0.5
    
    try:
        response = client.chat.completions.create(
            model=model,
            messages=[{
                "role": "user",
                "content": f"Summarize the following text into key points:\n{text}"
            }],
            temperature=temperature
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"❌ Error during text summarization: {str(e)}")
        return text  # Return original text if summarization fails

def generate_quiz(text, quiz_type, num_questions=5, difficulty="Standard"):
    # Check if we have a working client
    if not client_successful:
        return {
            "error": "No working AI provider available. Please configure an AI provider in the sidebar."
        }
    
    difficulty_instructions = {
        "Standard": """
        Create university-level questions that test comprehension, analysis, and application of the material.
        Questions should be straightforward but require good understanding of the content.
        Focus on key concepts, definitions, and logical connections.
        """,
        "Advanced": """
        Create advanced questions that require synthesis, evaluation, and critical thinking.
        Include scenario-based questions and complex problem-solving.
        Test ability to apply knowledge in new contexts.
        Questions should be challenging but fair, suitable for graduate-level study.
        """,
        "Extreme": """
        Create EXTREMELY challenging questions that require critical thinking, careful reading, and deep analysis.
        Make questions manipulative and tricky - use subtle distinctions, edge cases, and nuanced interpretations.
        Include questions that test ability to identify assumptions, logical fallacies, and hidden implications.
        Use complex scenarios that require synthesis of multiple concepts. Create sophisticated answer choices where:
        
        - Some questions may have multiple technically correct options, but only ONE is the BEST/MOST COMPLETE answer
        - Include "correct" vs "more correct" scenarios where students must choose the MOST accurate or comprehensive response
        - Design questions where 2-3 options could be partially right, but one stands out as superior
        
        Make incorrect options very plausible and tempting.
        Questions should be beyond university level - suitable for advanced professionals or doctoral students.
        """
    }
    
    # Get difficulty instruction or use Standard as fallback
    difficulty_instruction = difficulty_instructions.get(difficulty, difficulty_instructions["Standard"])
    
    if quiz_type == "Mixed (MCQ + T/F)":
        # For mixed quiz, create roughly half MCQ and half T/F
        mcq_count = num_questions // 2
        tf_count = num_questions - mcq_count
        
        prompt = f"""
Based on the following content, generate exactly {num_questions} questions:
- {mcq_count} multiple choice questions with 4 options each
- {tf_count} true or false questions

DIFFICULTY LEVEL: {difficulty}
{difficulty_instruction}

Return the response in this exact JSON format:
{{
  "questions": [
    {{
      "question": "Question text here",
      "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
      "correct_answer": "A",
      "explanation": "Explanation of why this is correct"
    }},
    {{
      "question": "Question text here",
      "options": ["True", "False"],
      "correct_answer": "True",
      "explanation": "Explanation of why this is correct"
    }}
  ]
}}

Content: {text}
"""
    elif quiz_type == "Multiple Choice":
        prompt = f"""
Based on the following content, generate exactly {num_questions} multiple choice questions with 4 options each.

DIFFICULTY LEVEL: {difficulty}
{difficulty_instruction}

Return the response in this exact JSON format:
{{
  "questions": [
    {{
      "question": "Question text here",
      "options": ["A) Option 1", "B) Option 2", "C) Option 3", "D) Option 4"],
      "correct_answer": "A",
      "explanation": "Explanation of why this is correct"
    }}
  ]
}}

Content: {text}
"""
    else:  # True or False
        prompt = f"""
Based on the following content, generate exactly {num_questions} true or false questions.

DIFFICULTY LEVEL: {difficulty}
{difficulty_instruction}

Return the response in this exact JSON format:
{{
  "questions": [
    {{
      "question": "Question text here",
      "options": ["True", "False"],
      "correct_answer": "True",
      "explanation": "Explanation of why this is correct"
    }}
  ]
}}

Content: {text}
"""
    
    # Use appropriate model based on provider
    if st.session_state.ai_provider == "Local AI (Ollama)":
        # Use selected model from session state if available, otherwise fall back to config
        model = getattr(st.session_state, 'selected_local_model', local_ai_config.MODEL_NAME)
        temperature = local_ai_config.TEMPERATURE
    elif st.session_state.ai_provider == "Google AI":
        model = google_ai_config.CHAT_MODEL
        temperature = google_ai_config.TEMPERATURE
    else:
        model = "gpt-3.5-turbo"
        temperature = 0.7
    
    try:
        response = client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=temperature
        )
        
        # Get the response content
        content = response.choices[0].message.content.strip()
        
        # Try to parse JSON directly first
        try:
            quiz_data = json.loads(content)
            return quiz_data
        except json.JSONDecodeError:
            pass
        
        # If direct parsing fails, try to extract JSON from markdown code blocks
        # Look for ```json ... ``` blocks - improved regex to capture complete JSON
        json_match = re.search(r'```(?:json)?\s*(\{.*\})\s*```', content, re.DOTALL)
        if json_match:
            try:
                quiz_data = json.loads(json_match.group(1))
                return quiz_data
            except json.JSONDecodeError as e:
                # Log the extraction attempt for debugging
                logger.warning(f"Failed to parse extracted JSON from code block: {e}")
                pass
        
        # Alternative: Look for JSON between triple backticks without json specifier
        json_match = re.search(r'```\s*(\{.*\})\s*```', content, re.DOTALL)
        if json_match:
            try:
                quiz_data = json.loads(json_match.group(1))
                return quiz_data
            except json.JSONDecodeError:
                pass
        
        # Look for any JSON-like structure in the entire content
        json_match = re.search(r'\{.*\}', content, re.DOTALL)
        if json_match:
            try:
                quiz_data = json.loads(json_match.group())
                return quiz_data
            except json.JSONDecodeError:
                pass
        
        # If all JSON parsing fails, return error with more context for debugging
        logger.error(f"Failed to parse quiz data. Full response: {content}")
        return {
            "error": f"Failed to parse quiz data. Response content: {content[:1000]}..." if len(content) > 1000 else f"Failed to parse quiz data. Response content: {content}"
        }
        
    except Exception as e:
        return {"error": f"Error processing response: {str(e)}"}

def display_quiz(quiz_data):
    """Display interactive quiz and handle user responses"""
    if "error" in quiz_data:
        st.error("Failed to generate quiz. Please try again.")
        return
    
    questions = quiz_data.get("questions", [])
    if not questions:
        st.error("No questions found in the quiz data.")
        return
    
    # Initialize session state for quiz
    if "current_question" not in st.session_state:
        st.session_state.current_question = 0
        st.session_state.user_answers = {}
        st.session_state.quiz_completed = False
        st.session_state.quiz_data = quiz_data
    
    # Reset quiz if new quiz is generated
    if st.session_state.quiz_data != quiz_data:
        st.session_state.current_question = 0
        st.session_state.user_answers = {}
        st.session_state.quiz_completed = False
        st.session_state.quiz_data = quiz_data
    
    total_questions = len(questions)
    
    if not st.session_state.quiz_completed:
        # Display current question
        current_q = st.session_state.current_question
        question_data = questions[current_q]
        
        # Progress bar
        progress = (current_q) / total_questions
        st.progress(progress)
        st.write(f"Question {current_q + 1} of {total_questions}")
        
        # Display question
        st.subheader(f"Q{current_q + 1}: {question_data['question']}")
        
        # Check if this is an open-ended question
        if question_data.get('type') == 'open_ended':
            # Display marking information
            total_marks = question_data.get('total_marks', 0)
            st.info(f"📝 **Open-ended Question** | Total Marks: {total_marks}")
            
            # Show word count guidance
            st.caption("💡 Write a comprehensive answer. Quality matters more than quantity!")
            
            # Text area for user input
            current_answer = st.session_state.user_answers.get(current_q, "")
            user_answer = st.text_area(
                "Your Answer:",
                value=current_answer,
                height=150,
                key=f"open_q_{current_q}",
                placeholder="Write your detailed answer here..."
            )
            
            # Word count display
            if user_answer:
                word_count = len(user_answer.split())
                st.caption(f"Word count: {word_count}")
            
        else:
            # Display options for MCQ/T/F questions
            options = question_data['options']
            if len(options) == 2:  # True/False
                user_answer = st.radio(
                    "Select your answer:",
                    options,
                    key=f"q_{current_q}",
                    index=None
                )
            else:  # Multiple choice
                user_answer = st.radio(
                    "Select your answer:",
                    options,
                    key=f"q_{current_q}",
                    index=None
                )
        
        col1, col2, col3 = st.columns([1, 1, 1])
        
        with col1:
            if current_q > 0:
                if st.button("Previous"):
                    st.session_state.current_question -= 1
                    st.rerun()
        
        with col3:
            # Check if user has provided an answer (different logic for open-ended vs MCQ/T/F)
            has_answer = False
            if question_data.get('type') == 'open_ended':
                has_answer = user_answer and user_answer.strip() and len(user_answer.split()) >= 5
            else:
                has_answer = user_answer is not None
            
            if has_answer:
                if current_q < total_questions - 1:
                    if st.button("Next"):
                        st.session_state.user_answers[current_q] = user_answer
                        st.session_state.current_question += 1
                        st.rerun()
                else:
                    if st.button("Submit Quiz"):
                        st.session_state.user_answers[current_q] = user_answer
                        st.session_state.quiz_completed = True
                        st.rerun()
            else:
                # Show guidance for what's needed
                if question_data.get('type') == 'open_ended':
                    st.caption("⚠️ Please write at least 5 words to proceed")
                else:
                    st.caption("⚠️ Please select an answer to proceed")
    
    else:
        # Display results
        display_results(questions, st.session_state.user_answers)

def display_results(questions, user_answers):
    """Display quiz results with detailed feedback"""
    st.success("🎉 Quiz Completed!")
    
    # Get difficulty and quiz type from session state
    difficulty = getattr(st.session_state, 'quiz_difficulty', 'Standard')
    quiz_type = getattr(st.session_state, 'quiz_type', 'Mixed')
    
    # Separate scoring for different question types
    traditional_questions = []
    open_ended_questions = []
    
    for i, question in enumerate(questions):
        if question.get('type') == 'open_ended':
            open_ended_questions.append((i, question))
        else:
            traditional_questions.append((i, question))
    
    # Calculate traditional question scores
    traditional_correct = 0
    total_traditional = len(traditional_questions)
    
    for i, question in traditional_questions:
        user_answer = user_answers.get(i, "")
        correct_answer = question['correct_answer']
        
        # For multiple choice, extract letter from user answer
        if len(question['options']) > 2:
            if user_answer and user_answer[0] in ['A', 'B', 'C', 'D']:
                user_letter = user_answer[0]
            else:
                user_letter = ""
        else:
            user_letter = user_answer
        
        if user_letter == correct_answer:
            traditional_correct += 1
    
    # Score open-ended questions
    open_ended_scores = []
    total_open_ended_marks = 0
    earned_open_ended_marks = 0
    
    if open_ended_questions and client_successful:
        st.info("🤖 Scoring open-ended questions with AI... This may take a moment.")
        progress_bar = st.progress(0)
        
        from open_ended_processor import OpenEndedQuestionProcessor
        processor = OpenEndedQuestionProcessor(
            client, 
            use_google_ai=(st.session_state.ai_provider == "Google AI"),
            use_local_ai=(st.session_state.ai_provider == "Local AI (Ollama)")
        )
        
        for idx, (i, question) in enumerate(open_ended_questions):
            user_answer = user_answers.get(i, "")
            scoring_result = processor.score_open_ended_answer(question, user_answer)
            open_ended_scores.append((i, question, scoring_result))
            
            total_open_ended_marks += scoring_result['max_score']
            earned_open_ended_marks += scoring_result['total_score']
            
            progress_bar.progress((idx + 1) / len(open_ended_questions))
        
        progress_bar.empty()
    
    # Calculate overall score
    if total_traditional > 0 and total_open_ended_marks > 0:
        # Mixed quiz - combine scores proportionally
        traditional_percentage = (traditional_correct / total_traditional) * 100
        open_ended_percentage = (earned_open_ended_marks / total_open_ended_marks) * 100
        
        # Weight equally for now (could be customized)
        overall_percentage = (traditional_percentage + open_ended_percentage) / 2
        
        st.subheader(f"📊 Overall Score: {overall_percentage:.1f}% {DIFFICULTY_CONFIG[difficulty]['emoji']} {difficulty} Level")
        
        # Show breakdown
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Traditional Questions", f"{traditional_correct}/{total_traditional} ({traditional_percentage:.1f}%)")
        with col2:
            st.metric("Open-ended Questions", f"{earned_open_ended_marks:.1f}/{total_open_ended_marks} ({open_ended_percentage:.1f}%)")
            
    elif total_open_ended_marks > 0:
        # Only open-ended questions
        overall_percentage = (earned_open_ended_marks / total_open_ended_marks) * 100
        st.subheader(f"📊 Overall Score: {earned_open_ended_marks:.1f}/{total_open_ended_marks} ({overall_percentage:.1f}%) {DIFFICULTY_CONFIG[difficulty]['emoji']} {difficulty} Level")
        
    else:
        # Only traditional questions
        overall_percentage = (traditional_correct / total_traditional) * 100
        st.subheader(f"📊 Overall Score: {traditional_correct}/{total_traditional} ({overall_percentage:.1f}%) {DIFFICULTY_CONFIG[difficulty]['emoji']} {difficulty} Level")
    
    # Score interpretation
    scoring_config = SCORING_CONFIG.get(difficulty, SCORING_CONFIG['Standard'])
    
    for level, (threshold, message) in scoring_config.items():
        if level == 'default':
            continue
        if overall_percentage >= threshold:
            if level == 'excellent':
                st.success(message)
            elif level == 'good':
                st.success(message)
            elif level == 'fair':
                st.info(message)
            else:
                st.warning(message)
            break
    else:
        st.warning(scoring_config['default'][1])
    
    # Detailed review
    st.subheader("📝 Detailed Review")
    
    for i, question in enumerate(questions):
        user_answer = user_answers.get(i, "No answer")
        
        if question.get('type') == 'open_ended':
            # Find the scoring result for this question
            scoring_result = None
            for scored_i, scored_q, result in open_ended_scores:
                if scored_i == i:
                    scoring_result = result
                    break
            
            if scoring_result:
                score_text = f"{scoring_result['total_score']:.1f}/{scoring_result['max_score']}"
                percentage = scoring_result['percentage']
                
                with st.expander(f"Question {i+1}: 📝 {score_text} ({percentage:.1f}%)"):
                    st.write(f"**Question:** {question['question']}")
                    st.write(f"**Your Answer:** {user_answer}")
                    st.write(f"**Score:** {score_text} marks ({percentage:.1f}%)")
                    
                    if scoring_result.get('overall_feedback'):
                        st.write(f"**Overall Feedback:** {scoring_result['overall_feedback']}")
                    
                    if scoring_result.get('criterion_scores'):
                        st.write("**Detailed Breakdown:**")
                        for criterion in scoring_result['criterion_scores']:
                            criterion_score = criterion['marks_awarded']
                            criterion_max = criterion['max_marks']
                            st.write(f"- {criterion['criterion']}: {criterion_score}/{criterion_max}")
                            if criterion.get('feedback'):
                                st.write(f"  *{criterion['feedback']}*")
                    
                    if scoring_result.get('strengths'):
                        st.success("**Strengths:** " + ", ".join(scoring_result['strengths']))
                    
                    if scoring_result.get('improvements'):
                        st.info("**Areas for Improvement:** " + ", ".join(scoring_result['improvements']))
                    
                    # Show model answer
                    with st.expander("View Model Answer"):
                        st.write(question['model_answer'])
            
        else:
            # Traditional MCQ/T/F questions
            correct_answer = question['correct_answer']
            
            # For multiple choice, extract letter from user answer
            if len(question['options']) > 2:
                if user_answer and len(user_answer) > 0 and user_answer[0] in ['A', 'B', 'C', 'D']:
                    user_letter = user_answer[0]
                else:
                    user_letter = ""
            else:
                user_letter = user_answer
            
            is_correct = user_letter == correct_answer
            
            with st.expander(f"Question {i+1}: {'✅' if is_correct else '❌'}"):
                st.write(f"**Question:** {question['question']}")
                st.write(f"**Your Answer:** {user_answer}")
                
                # Find correct answer text
                if len(question['options']) > 2:
                    correct_option = next((opt for opt in question['options'] if opt[0] == correct_answer), f"{correct_answer})")
                else:
                    correct_option = correct_answer
                    
                st.write(f"**Correct Answer:** {correct_option}")
                st.write(f"**Explanation:** {question['explanation']}")
                
                if is_correct:
                    st.success("Correct! 🎉")
                else:
                    st.error("Incorrect 😔")
    
    # Restart quiz button
    if st.button("Take Quiz Again"):
        st.session_state.current_question = 0
        st.session_state.user_answers = {}
        st.session_state.quiz_completed = False
        st.rerun()

def display_study_materials(materials_data, material_type):
    """Display generated study materials with appropriate formatting."""
    
    if "error" in materials_data:
        st.error(f"❌ Failed to generate {material_type.lower()}: {materials_data['error']}")
        return
    
    st.success(f"✅ {material_type} generated successfully!")
    
    if material_type == "Complete Study Guide":
        display_complete_study_guide(materials_data)
    elif material_type == "Summary Only":
        display_summary(materials_data)
    elif material_type == "Cheat Sheet":
        display_cheat_sheet(materials_data)
    elif material_type == "Flashcards":
        display_flashcards(materials_data)
    elif material_type == "Study Outline":
        display_outline(materials_data)
    elif material_type == "Key Terms":
        display_key_terms(materials_data)

def display_complete_study_guide(guide_data):
    """Display a complete study guide with all components."""
    
    st.subheader(f"📚 {guide_data.get('title', 'Study Guide')}")
    st.caption(f"Generated: {guide_data.get('generated_at', 'Unknown time')} | Type: {guide_data.get('guide_type', 'Unknown').title()}")
    
    # Show study plan first
    if 'study_plan' in guide_data:
        with st.expander("🗓️ Suggested Study Plan", expanded=True):
            plan = guide_data['study_plan']
            st.info(f"**Total Study Time:** {plan.get('total_time', 'Variable')}")
            
            for session in plan.get('sessions', []):
                st.write(f"**Session {session['session']}** ({session['time']}): {session['focus']}")
    
    # Display each component
    components = guide_data.get('components', {})
    
    # Summary
    if 'summary' in components and components['summary']:
        with st.expander("📖 Summary", expanded=True):
            display_summary(components['summary'])
    
    # Key Terms
    if 'key_terms' in components and components['key_terms']:
        with st.expander("📚 Key Terms & Definitions", expanded=False):
            display_key_terms(components['key_terms'])
    
    # Cheat Sheet
    if 'cheat_sheet' in components and components['cheat_sheet']:
        with st.expander("📄 Quick Reference Cheat Sheet", expanded=False):
            display_cheat_sheet(components['cheat_sheet'])
    
    # Flashcards
    if 'flashcards' in components and components['flashcards']:
        with st.expander("🔄 Interactive Flashcards", expanded=False):
            display_flashcards(components['flashcards'])
    
    # Show any errors
    if guide_data.get('errors'):
        with st.expander("⚠️ Generation Notes"):
            for error in guide_data['errors']:
                st.warning(error)

def display_summary(summary_data):
    """Display a summary with key points and topics."""
    
    if "error" in summary_data:
        st.error(f"❌ {summary_data['error']}")
        return
    
    # Main summary
    summary_text = summary_data.get('summary', 'No summary available')
    st.write(summary_text)
    
    # Metadata
    col1, col2 = st.columns(2)
    with col1:
        if 'word_count' in summary_data:
            st.metric("Word Count", summary_data['word_count'])
    with col2:
        if 'summary_type' in summary_data:
            st.metric("Type", summary_data['summary_type'].title())
    
    # Key points
    key_points = summary_data.get('key_points', [])
    if key_points:
        st.subheader("🎯 Key Points")
        for point in key_points:
            st.write(f"• {point}")
    
    # Main topics
    main_topics = summary_data.get('main_topics', [])
    if main_topics:
        st.subheader("📋 Main Topics")
        topic_cols = st.columns(3)
        for i, topic in enumerate(main_topics):
            with topic_cols[i % 3]:
                st.write(f"📌 {topic}")

def display_cheat_sheet(cheat_data):
    """Display a formatted cheat sheet."""
    
    if "error" in cheat_data:
        st.error(f"❌ {cheat_data['error']}")
        return
    
    st.subheader(f"📋 {cheat_data.get('title', 'Cheat Sheet')}")
    
    # Sections
    sections = cheat_data.get('sections', [])
    for section in sections:
        st.subheader(f"📌 {section.get('heading', 'Section')}")
        st.write(section.get('content', ''))
        
        items = section.get('items', [])
        if items:
            for item in items:
                st.write(f"• {item}")
    
    # Key terms in a nice layout
    key_terms = cheat_data.get('key_terms', [])
    if key_terms:
        st.subheader("📚 Key Terms")
        for term_data in key_terms:
            with st.container():
                st.write(f"**{term_data.get('term', 'Term')}**: {term_data.get('definition', 'Definition')}")
    
    # Formulas
    formulas = cheat_data.get('formulas', [])
    if formulas:
        st.subheader("🔢 Formulas")
        for formula in formulas:
            st.write(f"**{formula.get('name', 'Formula')}**: `{formula.get('formula', 'N/A')}`")
            if formula.get('explanation'):
                st.caption(formula['explanation'])
    
    # Quick tips
    tips = cheat_data.get('quick_tips', [])
    if tips:
        st.subheader("💡 Quick Tips")
        for tip in tips:
            st.info(tip)

def display_flashcards(flashcard_data):
    """Display interactive flashcards."""
    
    if "error" in flashcard_data:
        st.error(f"❌ {flashcard_data['error']}")
        return
    
    flashcards = flashcard_data.get('flashcards', [])
    if not flashcards:
        st.warning("No flashcards generated.")
        return
    
    # Initialize flashcard session state
    if "current_flashcard" not in st.session_state:
        st.session_state.current_flashcard = 0
        st.session_state.flashcard_answer_visible = False
        st.session_state.flashcard_stats = {"correct": 0, "incorrect": 0, "skipped": 0}
    
    current_card = flashcards[st.session_state.current_flashcard]
    total_cards = len(flashcards)
    
    # Progress and stats
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Progress", f"{st.session_state.current_flashcard + 1}/{total_cards}")
    with col2:
        st.metric("Correct", st.session_state.flashcard_stats["correct"])
    with col3:
        st.metric("Difficulty", current_card.get('difficulty', 'unknown').title())
    
    # Card display
    with st.container():
        st.subheader(f"🔄 Card {st.session_state.current_flashcard + 1}")
        
        # Show category and difficulty
        category = current_card.get('category', 'General')
        difficulty = current_card.get('difficulty', 'unknown')
        st.caption(f"Category: {category} | Difficulty: {difficulty.title()}")
        
        # Front of card
        st.markdown("### 📝 Question:")
        st.write(current_card.get('front', 'No question available'))
        
        # Hint if available
        if current_card.get('hint') and not st.session_state.flashcard_answer_visible:
            with st.expander("💡 Hint"):
                st.write(current_card['hint'])
        
        # Answer section
        if st.session_state.flashcard_answer_visible:
            st.markdown("### ✅ Answer:")
            st.success(current_card.get('back', 'No answer available'))
            
            # Self-assessment buttons
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                if st.button("😊 Got it right!", key="correct"):
                    st.session_state.flashcard_stats["correct"] += 1
                    next_card()
            with col2:
                if st.button("😔 Got it wrong", key="incorrect"):
                    st.session_state.flashcard_stats["incorrect"] += 1
                    next_card()
            with col3:
                if st.button("⏭️ Skip", key="skip"):
                    st.session_state.flashcard_stats["skipped"] += 1
                    next_card()
            with col4:
                if st.button("🔄 Flip Back", key="flip_back"):
                    st.session_state.flashcard_answer_visible = False
                    st.rerun()
        else:
            if st.button("🔄 Show Answer", key="show_answer", type="primary"):
                st.session_state.flashcard_answer_visible = True
                st.rerun()
    
    # Navigation
    col1, col2, col3 = st.columns(3)
    with col1:
        if st.button("⬅️ Previous") and st.session_state.current_flashcard > 0:
            st.session_state.current_flashcard -= 1
            st.session_state.flashcard_answer_visible = False
            st.rerun()
    with col2:
        if st.button("🔄 Shuffle Cards"):
            import random
            random.shuffle(flashcards)
            flashcard_data['flashcards'] = flashcards
            st.session_state.current_flashcard = 0
            st.session_state.flashcard_answer_visible = False
            st.success("Cards shuffled!")
            st.rerun()
    with col3:
        if st.button("➡️ Next") and st.session_state.current_flashcard < total_cards - 1:
            st.session_state.current_flashcard += 1
            st.session_state.flashcard_answer_visible = False
            st.rerun()
    
    # Study tips
    study_tips = flashcard_data.get('study_tips', [])
    if study_tips:
        with st.expander("💡 Study Tips"):
            for tip in study_tips:
                st.write(f"• {tip}")

def next_card():
    """Move to next flashcard."""
    if st.session_state.current_flashcard < len(st.session_state.get('flashcards', [])) - 1:
        st.session_state.current_flashcard += 1
    else:
        st.session_state.current_flashcard = 0  # Loop back to beginning
    st.session_state.flashcard_answer_visible = False
    st.rerun()

def display_outline(outline_data):
    """Display a structured study outline."""
    
    if "error" in outline_data:
        st.error(f"❌ {outline_data['error']}")
        return
    
    # Metadata
    col1, col2, col3 = st.columns(3)
    with col1:
        total_sections = outline_data.get('total_sections', 0)
        st.metric("Sections", total_sections)
    with col2:
        max_depth = outline_data.get('max_depth', 0)
        st.metric("Max Depth", max_depth)
    with col3:
        total_time = outline_data.get('time_estimates', {}).get('total_study_time', 'Unknown')
        st.metric("Study Time", total_time)
    
    # Outline structure
    outline = outline_data.get('outline', [])
    render_outline_recursive(outline, 0)
    
    # Study sequence
    sequence = outline_data.get('study_sequence', [])
    if sequence:
        st.subheader("📅 Recommended Study Sequence")
        for i, section in enumerate(sequence, 1):
            time_estimate = outline_data.get('time_estimates', {}).get('per_section', ['30 min'] * len(sequence))
            time_for_section = time_estimate[i-1] if i-1 < len(time_estimate) else '30 min'
            st.write(f"{i}. **{section}** ({time_for_section})")

def render_outline_recursive(outline_items, depth):
    """Recursively render outline structure."""
    for item in outline_items:
        level = item.get('level', 1)
        marker = item.get('marker', '')
        text = item.get('text', '')
        
        # Indent based on level
        indent = "  " * (level - 1)
        if level == 1:
            st.subheader(f"{marker}. {text}")
        elif level == 2:
            st.write(f"**{indent}{marker}. {text}**")
        else:
            st.write(f"{indent}{marker}. {text}")
        
        # Recursively render children
        children = item.get('children', [])
        if children:
            render_outline_recursive(children, depth + 1)

def display_key_terms(terms_data):
    """Display key terms and definitions."""
    
    if "error" in terms_data:
        st.error(f"❌ {terms_data['error']}")
        return
    
    # Metadata
    total_terms = terms_data.get('total_terms', 0)
    st.metric("Total Terms", total_terms)
    
    # Categories
    categories = terms_data.get('categories', [])
    if categories:
        st.subheader("📂 Categories")
        for category in categories:
            st.write(f"**{category.get('category', 'Category')}**: {len(category.get('terms', []))} terms")
    
    # Terms list
    key_terms = terms_data.get('key_terms', [])
    if key_terms:
        st.subheader("📚 Terms & Definitions")
        
        # Group by importance
        high_importance = [term for term in key_terms if term.get('importance') == 'high']
        medium_importance = [term for term in key_terms if term.get('importance') == 'medium']
        low_importance = [term for term in key_terms if term.get('importance') == 'low']
        
        if high_importance:
            st.markdown("### 🔴 High Priority Terms")
            for term in high_importance:
                display_term(term)
        
        if medium_importance:
            st.markdown("### 🟡 Medium Priority Terms")
            for term in medium_importance:
                display_term(term)
                
        if low_importance:
            st.markdown("### 🟢 Low Priority Terms")
            for term in low_importance:
                display_term(term)
    
    # Study suggestions
    suggestions = terms_data.get('study_suggestions', [])
    if suggestions:
        st.subheader("💡 Study Suggestions")
        for suggestion in suggestions:
            st.info(suggestion)

def display_term(term):
    """Display a single term with its definition and context."""
    with st.expander(f"📖 {term.get('term', 'Term')}"):
        st.write(f"**Definition**: {term.get('definition', 'No definition available')}")
        
        context = term.get('context', '')
        if context:
            st.write(f"**Context**: {context}")
        
        related_terms = term.get('related_terms', [])
        if related_terms:
            st.write(f"**Related Terms**: {', '.join(related_terms)}")

def generate_quiz_content(final_text, quiz_type, num_questions, difficulty, mcq_count=0, tf_count=0, open_count=0):
    """Generate quiz content with error handling."""
    with st.spinner(f"🤖 Generating quiz using {ai_provider}..."):
        try:
            # Handle different quiz types
            if quiz_type == "Open-ended Questions":
                from open_ended_processor import OpenEndedQuestionProcessor
                processor = OpenEndedQuestionProcessor(
                    client, 
                    use_google_ai=(st.session_state.ai_provider == "Google AI"),
                    use_local_ai=(st.session_state.ai_provider == "Local AI (Ollama)")
                )
                quiz_data = processor.generate_open_ended_questions(final_text, num_questions, difficulty)
            elif quiz_type == "Complete Mix (All Types)":
                from open_ended_processor import OpenEndedQuestionProcessor
                processor = OpenEndedQuestionProcessor(
                    client, 
                    use_google_ai=(st.session_state.ai_provider == "Google AI"),
                    use_local_ai=(st.session_state.ai_provider == "Local AI (Ollama)")
                )
                quiz_data = processor.generate_mixed_quiz(final_text, mcq_count, tf_count, open_count, difficulty)
            else:
                quiz_data = generate_quiz(final_text, quiz_type, num_questions, difficulty)
            
            if "error" not in quiz_data:
                st.session_state.quiz_generated = True
                st.session_state.quiz_data = quiz_data
                st.session_state.quiz_difficulty = difficulty  # Store difficulty for results
                st.session_state.quiz_type = quiz_type  # Store quiz type
                st.success("✅ Quiz generated successfully! Start answering below.")
                st.rerun()
            else:
                st.error(f"❌ Failed to generate quiz: {quiz_data.get('error', 'Unknown error')}")
                
        except Exception as e:
            st.error(f"❌ Error during quiz generation: {str(e)}")
            # Add debug information
            st.write("**Debug Info:**")
            st.write(f"- AI Provider: {ai_provider}")
            st.write(f"- Quiz Type: {quiz_type}")
            st.write(f"- Text Length: {len(final_text)} characters")
            if app_config.DEBUG_MODE:
                st.exception(e)

def generate_study_materials_content(final_text, material_type, local_vars):
    """Generate study materials content with error handling."""
    with st.spinner(f"📚 Generating {material_type.lower()} using {ai_provider}..."):
        try:
            # Initialize study materials generator
            materials_generator = StudyMaterialsGenerator(
                client,
                use_google_ai=(st.session_state.ai_provider == "Google AI"),
                use_local_ai=(st.session_state.ai_provider == "Local AI (Ollama)")
            ) 
            
            # Generate based on material type
            if material_type == "Complete Study Guide":
                guide_type = local_vars.get('guide_type', 'comprehensive')
                materials_data = materials_generator.generate_study_guide(final_text, guide_type)
            elif material_type == "Summary Only":
                summary_type = local_vars.get('summary_type', 'detailed')
                materials_data = materials_generator.generate_comprehensive_summary(final_text, summary_type)
            elif material_type == "Cheat Sheet":
                cheat_format = local_vars.get('cheat_format', 'comprehensive')
                materials_data = materials_generator.generate_cheat_sheet(final_text, cheat_format)
            elif material_type == "Flashcards":
                card_count = local_vars.get('card_count', 15)
                flashcard_difficulty = local_vars.get('flashcard_difficulty', 'mixed')
                materials_data = materials_generator.generate_flashcards(final_text, card_count, flashcard_difficulty)
            elif material_type == "Study Outline":
                outline_depth = local_vars.get('outline_depth', 'detailed')
                materials_data = materials_generator.generate_study_outline(final_text, outline_depth)
            elif material_type == "Key Terms":
                term_count = local_vars.get('term_count', 15)
                materials_data = materials_generator.generate_key_terms(final_text, term_count)
            else:
                materials_data = {"error": f"Unknown material type: {material_type}"}
            
            if "error" not in materials_data:
                st.session_state.materials_generated = True
                st.session_state.materials_data = materials_data
                st.session_state.material_type = material_type
                st.success(f"✅ {material_type} generated successfully!")
                st.rerun()
            else:
                st.error(f"❌ Failed to generate {material_type.lower()}: {materials_data.get('error', 'Unknown error')}")
                
        except Exception as e:
            st.error(f"❌ Error during {material_type.lower()} generation: {str(e)}")
            # Add debug information
            st.write("**Debug Info:**")
            st.write(f"- AI Provider: {ai_provider}")
            st.write(f"- Material Type: {material_type}")
            st.write(f"- Text Length: {len(final_text)} characters")
            if app_config.DEBUG_MODE:
                st.exception(e)

def main():
    global client, ai_provider, client_successful
    
    st.title("📚 AI Interactive Quiz & Study Materials Generator")
    
    # Display AI provider info - use the actual working provider, not the selected one
    if ai_provider == "Local AI (Ollama)":
        provider_emoji = "🏠"
        provider_color = "orange" 
    elif ai_provider == "Google AI":
        provider_emoji = "🆕"
        provider_color = "green"
    else:
        provider_emoji = "⚡"
        provider_color = "blue"
    
    st.info(f"{provider_emoji} **Powered by {ai_provider}** - Advanced AI for intelligent quiz generation and study materials creation")
    
    # Initialize session state
    if "quiz_generated" not in st.session_state:
        st.session_state.quiz_generated = False
    if "materials_generated" not in st.session_state:
        st.session_state.materials_generated = False
    if "text_summarized" not in st.session_state:
        st.session_state.text_summarized = False
    if "summarized_text" not in st.session_state:
        st.session_state.summarized_text = ""
    if "original_text" not in st.session_state:
        st.session_state.original_text = ""
    if "summarization_in_progress" not in st.session_state:
        st.session_state.summarization_in_progress = False
    
    # Sidebar for quiz configuration
    with st.sidebar:
        st.header("Quiz Configuration")
        
        # AI Provider Selection with status
        selected_provider = session_manager.render_provider_selector()
        
        # If provider changed, reinitialize client
        if selected_provider != st.session_state.ai_provider:
            # Update global variables
            st.session_state.ai_provider = selected_provider
            client, ai_provider, client_successful = ai_factory.get_working_client()
            st.rerun()
        
        # API Key Configuration
        session_manager.render_api_key_inputs()
        
        st.markdown("---")
        
        uploaded_file = st.file_uploader("Upload PDF, Word, or PPTX file", type=["pdf", "docx", "pptx"])
        
        # Content generation type selector
        generation_type = st.selectbox(
            "Choose Generation Type",
            ["Interactive Quiz", "Study Materials"],
            help="Select whether to generate a quiz or study materials"
        )
        
        if generation_type == "Interactive Quiz":
            quiz_type = st.selectbox(
                "Choose Quiz Type", 
                ["Multiple Choice", "True or False", "Mixed (MCQ + T/F)", "Open-ended Questions", "Complete Mix (All Types)"]
            )
            difficulty = st.selectbox(
                "Choose Difficulty Level", 
                ["Standard", "Advanced", "Extreme"],
                index=0,  # Default to Standard
                help="Standard: University-level | Advanced: Graduate-level | Extreme: Expert-level with tricky elements"
            )
            
            # Adjust question count based on quiz type
            if quiz_type == "Open-ended Questions":
                num_questions = st.slider("Number of Questions", min_value=2, max_value=5, value=3)
                if st.session_state.ai_provider not in ["Local AI (Ollama)", "Google AI"]:
                    st.warning("💡 Open-ended questions use GPT-4 for scoring and may increase API costs. Each answer requires an additional AI evaluation.")
            elif quiz_type == "Complete Mix (All Types)":
                st.write("**Question Distribution:**")
                mcq_count = st.slider("Multiple Choice", min_value=1, max_value=5, value=2)
                tf_count = st.slider("True/False", min_value=1, max_value=5, value=2) 
                open_count = st.slider("Open-ended", min_value=1, max_value=3, value=1)
                num_questions = mcq_count + tf_count + open_count
                st.info(f"Total questions: {num_questions}")
                if open_count > 0 and st.session_state.ai_provider not in ["Local AI (Ollama)", "Google AI"]:
                    st.warning(f"⚠️ {open_count} open-ended question(s) will use GPT-4 for scoring (higher cost)")
            else:
                num_questions = st.slider("Number of Questions", min_value=3, max_value=10, value=5)
        
        else:  # Study Materials
            material_type = st.selectbox(
                "Choose Study Material Type",
                ["Complete Study Guide", "Summary Only", "Cheat Sheet", "Flashcards", "Study Outline", "Key Terms"],
                help="Select the type of study material to generate"
            )
            
            if material_type == "Complete Study Guide":
                guide_type = st.selectbox(
                    "Study Guide Type",
                    ["comprehensive", "exam_prep", "quick_review"],
                    format_func=lambda x: {
                        "comprehensive": "📚 Comprehensive Guide (4-6 hours study time)",
                        "exam_prep": "🎯 Exam Preparation (6-8 hours study time)", 
                        "quick_review": "⚡ Quick Review (2-3 hours study time)"
                    }[x]
                )
            elif material_type == "Summary Only":
                summary_type = st.selectbox(
                    "Summary Type",
                    ["detailed", "concise", "bullet_points"],
                    format_func=lambda x: {
                        "detailed": "📖 Detailed Summary (300-500 words)",
                        "concise": "📝 Concise Summary (150-250 words)",
                        "bullet_points": "• Bullet Points Summary"
                    }[x]
                )
            elif material_type == "Cheat Sheet":
                cheat_format = st.selectbox(
                    "Cheat Sheet Format",
                    ["comprehensive", "formulas", "definitions", "quick_ref"],
                    format_func=lambda x: {
                        "comprehensive": "📋 Comprehensive Reference",
                        "formulas": "🔢 Formulas & Equations",
                        "definitions": "📚 Definitions & Terms",
                        "quick_ref": "⚡ Quick Reference"
                    }[x]
                )
            elif material_type == "Flashcards":
                card_count = st.slider("Number of Flashcards", min_value=5, max_value=30, value=15)
                flashcard_difficulty = st.selectbox(
                    "Flashcard Difficulty",
                    ["basic", "intermediate", "advanced", "mixed"],
                    index=3,  # Default to mixed
                    format_func=lambda x: {
                        "basic": "📚 Basic (Definitions & Facts)",
                        "intermediate": "🎓 Intermediate (Concepts & Applications)",
                        "advanced": "🏆 Advanced (Analysis & Synthesis)",
                        "mixed": "🎯 Mixed Difficulty"
                    }[x]
                )
            elif material_type == "Study Outline":
                outline_depth = st.selectbox(
                    "Outline Depth",
                    ["overview", "detailed", "comprehensive"],
                    index=1,  # Default to detailed
                    format_func=lambda x: {
                        "overview": "📝 Overview (1-2 levels)",
                        "detailed": "📋 Detailed (3-4 levels)", 
                        "comprehensive": "📚 Comprehensive (4-5 levels)"
                    }[x]
                )
            elif material_type == "Key Terms":
                term_count = st.slider("Number of Key Terms", min_value=5, max_value=30, value=15)
        
        if uploaded_file:
            st.success("✅ File uploaded successfully!")
        
        # Show local AI status if selected
        if st.session_state.ai_provider == "Local AI (Ollama)":
            st.markdown("---")
            st.subheader("🏠 Local AI Status")
            ollama_base_url = local_ai_config.BASE_URL.replace('/v1', '')
            if is_ollama_running(ollama_base_url):
                st.success("✅ Ollama server running")
                available_models = list_available_models(ollama_base_url)
                if available_models:
                    # Initialize selected model in session state
                    if "selected_local_model" not in st.session_state:
                        # Default to current config model if available, otherwise first available
                        default_model = local_ai_config.MODEL_NAME if local_ai_config.MODEL_NAME in available_models else available_models[0]
                        st.session_state.selected_local_model = default_model
                    
                    # Model selector with performance info
                    model_options = []
                    for model in available_models:
                        model_options.append(model)
                    
                    selected_model = st.selectbox(
                        "🤖 Select Model:",
                        model_options,
                        format_func=lambda x: f"{get_model_info(x)}",
                        index=model_options.index(st.session_state.selected_local_model) if st.session_state.selected_local_model in model_options else 0,
                        key="model_selector",
                        help="Choose which model to use for quiz generation. Larger models are more capable but slower."
                    )
                    
                    # Update session state and config if model changed
                    if selected_model != st.session_state.selected_local_model:
                        st.session_state.selected_local_model = selected_model
                        # Update the local config to use the selected model
                        local_ai_config.MODEL_NAME = selected_model
                        st.success(f"🔄 Switched to model: {selected_model}")
                        
                        # Show performance hint for selected model
                        if ":2b" in selected_model:
                            st.info("⚡ **Fast & Efficient** - Good for quick quiz generation")
                        elif ":9b" in selected_model:
                            st.info("⚖️ **Balanced** - Good mix of speed and quality")
                        elif ":27b" in selected_model:
                            st.info("🎯 **High Quality** - Better responses, requires more time")
                        elif ":70b" in selected_model:
                            st.info("🏆 **Premium Quality** - Best results, much slower")
                        
                        st.rerun()
                    
                    # Show current model info
                    model_info = get_model_info(st.session_state.selected_local_model)
                    st.info(f"🎯 **Active Model:** {st.session_state.selected_local_model} {model_info}")
                    
                    # Show all available models in expander
                    with st.expander(f"📦 All Available Models ({len(available_models)})"):
                        for i, model in enumerate(available_models, 1):
                            is_current = model == st.session_state.selected_local_model
                            marker = "🔹 **" if is_current else "• "
                            end_marker = "** (Active)" if is_current else ""
                            model_info = get_model_info(model)
                            st.write(f"{marker}{model} {model_info}{end_marker}")
                else:
                    st.warning("No models found. Run `ollama pull gemma2:2b`")
            else:
                st.error("❌ Ollama server not running")
                st.code("ollama serve")

    if uploaded_file and not (st.session_state.quiz_generated or st.session_state.materials_generated):
        ext = uploaded_file.name.split(".")[-1]
        
        # Create a unique identifier for the uploaded file
        file_id = f"{uploaded_file.name}_{uploaded_file.size}"
        current_file_id = getattr(st.session_state, 'current_file_id', None)
        
        # Check if this is a new file
        if file_id != current_file_id:
            # Reset all states for new file
            st.session_state.current_file_id = file_id
            st.session_state.text_summarized = False
            st.session_state.summarized_text = ""
            st.session_state.original_text = ""
            st.session_state.summarization_in_progress = False
            st.session_state.quiz_generated = False
        
        # Extract text only if we don't have it already
        if not st.session_state.original_text:
            try:
                if ext == "pdf":
                    text = extract_text_from_pdf(uploaded_file)
                elif ext == "docx":
                    text = extract_text_from_docx(uploaded_file)
                elif ext == "pptx":
                    text = extract_text_from_pptx(uploaded_file)
                else:
                    st.error("❌ Unsupported file format")
                    return
            except Exception as e:
                st.error(f"Error reading file: {str(e)}")
                return

            if not text.strip():
                st.error("❌ No text found in the uploaded file")
                return
                
            # Store the original text
            st.session_state.original_text = text
        else:
            # Use the stored text
            text = st.session_state.original_text

        # Show document preview
        with st.expander("📄 Document Preview"):
            preview_text = st.session_state.summarized_text if st.session_state.text_summarized else text
            st.text_area("Extracted Text", preview_text[:1000] + "..." if len(preview_text) > 1000 else preview_text, height=200)

        # Handle summarization logic
        needs_summarization = len(text) > 3000 and not st.session_state.text_summarized
        
        if needs_summarization and not st.session_state.summarization_in_progress:
            # Start summarization automatically
            st.session_state.summarization_in_progress = True
            st.info("📄 Large content detected. Summarizing automatically...")
            
            # Check if we have a working client for summarization
            if not client_successful:
                st.warning("⚠️ No AI provider available for summarization. Using original text.")
                st.session_state.summarized_text = text
                st.session_state.text_summarized = True
                st.session_state.summarization_in_progress = False
                st.rerun()
            else:
                with st.spinner("Summarizing content..."):
                    summarized = summarize_text(text)
                    st.session_state.summarized_text = summarized
                    st.session_state.text_summarized = True
                    st.session_state.summarization_in_progress = False
                    st.success("✅ Content summarized successfully!")
                    st.rerun()
        
        # Show summarization status
        if st.session_state.summarization_in_progress:
            st.info("🔄 Summarization in progress... Please wait.")
            st.stop()  # Prevent the rest of the UI from rendering
        elif st.session_state.text_summarized:
            st.success(f"✅ Content summarized (Original: {len(text):,} chars → Summary: {len(st.session_state.summarized_text):,} chars)")
            
        # Determine which text to use for quiz generation
        final_text = st.session_state.summarized_text if st.session_state.text_summarized else text
        
        # Generation button - only show if summarization is complete (if needed)
        if needs_summarization and not st.session_state.text_summarized:
            st.info("⏳ Please wait for summarization to complete before generating content.")
        else:
            if generation_type == "Interactive Quiz":
                button_text = "🎯 Generate Interactive Quiz"
                button_help = "Create an interactive quiz from your document"
            else:
                button_text = f"📚 Generate {material_type}"
                button_help = f"Create {material_type.lower()} from your document"
                
            if st.button(button_text, type="primary", help=button_help):
                # Check if we have a working client before attempting generation
                if not client_successful:
                    st.error("❌ No working AI provider available. Please configure an AI provider in the sidebar first.")
                    st.info("💡 **Quick Setup Guide:**")
                    st.info("1. **Local AI**: Start Ollama server (`ollama serve`) and pull a model (`ollama pull gemma2:2b`)")
                    st.info("2. **Google AI**: Enter your Google AI API key in the sidebar")
                    st.info("3. **OpenAI**: Enter your OpenAI API key in the sidebar")
                    return
                
                if generation_type == "Interactive Quiz":
                    generate_quiz_content(final_text, quiz_type, num_questions, difficulty, mcq_count if 'mcq_count' in locals() else 0, tf_count if 'tf_count' in locals() else 0, open_count if 'open_count' in locals() else 0)
                else:
                    generate_study_materials_content(final_text, material_type, locals())

    elif uploaded_file and st.session_state.quiz_generated:
        # Display the interactive quiz
        st.markdown("---")
        display_quiz(st.session_state.quiz_data)
        
        # Reset quiz button in sidebar
        with st.sidebar:
            st.markdown("---")
            if st.button("🔄 Generate New Quiz"):
                st.session_state.quiz_generated = False
                st.session_state.current_question = 0
                st.session_state.user_answers = {}
                st.session_state.quiz_completed = False
                # Keep the summarized text to avoid re-summarization
                st.rerun()
    
    elif uploaded_file and st.session_state.materials_generated:
        # Display the study materials
        st.markdown("---")
        display_study_materials(st.session_state.materials_data, st.session_state.material_type)
        
        # Reset materials button in sidebar
        with st.sidebar:
            st.markdown("---")
            if st.button("🔄 Generate New Materials"):
                st.session_state.materials_generated = False
                # Keep the summarized text to avoid re-summarization
                st.rerun()
    
    else:
        # Welcome screen with improved guidance
        st.markdown("""
        ## Welcome to the AI Interactive Quiz Generator & Study Materials Creator! 🎓
        
        This application creates personalized, interactive quizzes and comprehensive study materials from your documents with **advanced AI scoring** and **graceful error handling**.
        
        ### 🆕 **New Features:**
        - **📝 Study Materials Generation**: Create summaries, cheat sheets, flashcards, and more!
        - **🔧 Graceful Error Handling**: App now works even if AI providers have issues
        - **⚡ Dynamic Provider Switching**: Change AI providers instantly without restarting
        - **🔐 Runtime API Key Management**: Enter API keys directly in the app and save them
        - **💾 Configuration Persistence**: Optionally save your settings for next time
        
        ### How it works:
        1. 📁 **Upload** a PDF, Word, or PowerPoint file
        2. ⚙️ **Configure** your AI provider and API keys in the sidebar
        3. 🎯 **Choose Generation Type**: Interactive Quiz or Study Materials
        4. 📝 **Generate & Use** your personalized content
        5. 📊 **Review** results with detailed AI-powered feedback
        
        ### AI Provider Options:
        - 🏠 **Local AI (Ollama)**: Run Gemma models locally - completely free and private!
        - 🆕 **Google AI**: Use Google's Gemma models via API
        - ⚡ **OpenAI**: Traditional GPT models for reference
        
        ### 📚 **Study Materials Available:**
        - **📖 Complete Study Guide**: Comprehensive package with summary, cheat sheet, flashcards, key terms + study plan
        - **📝 Summaries**: Detailed, concise, or bullet-point summaries
        - **📋 Cheat Sheets**: Quick reference sheets with key concepts
        - **🔄 Interactive Flashcards**: Self-paced flashcards with difficulty levels
        - **📊 Study Outlines**: Structured hierarchical outlines (available separately)
        - **🔖 Key Terms**: Important terminology with definitions
        
        ### Question Types Available:
        - 🔘 **Multiple Choice**: Traditional 4-option questions
        - ✅ **True/False**: Binary choice questions  
        - 📝 **Open-ended**: Write detailed answers scored by AI
        - 🎯 **Complete Mix**: Combination of all question types
        
        ### 🔐 **New API Key Management:**
        - Enter API keys directly in the sidebar
        - Choose to save them locally for next session
        - No more .env file editing required!
        - Secure password input fields
        
        **Get started by:**
        1. **Configure AI Provider** in the sidebar →
        2. **Enter API keys** (if needed)
        3. **Upload a document** to begin!
        """)
        
        # Provider status overview
        with st.expander("🔍 Current Provider Status"):
            session_manager.update_provider_status()
            for provider, status in st.session_state.provider_status.items():
                if status.get("available", False):
                    st.success(f"✅ **{provider}**: {status.get('message', 'Ready')}")
                else:
                    st.warning(f"⚠️ **{provider}**: {status.get('message', 'Not available')}")
        
        # Study Materials Preview
        with st.expander("📚 Study Materials Preview"):
            st.markdown("""
            **Example Study Materials from a Biology Document:**
            
            **📖 Summary**: "Cell membrane is a selectively permeable barrier that controls substance transport..."
            
            **📋 Cheat Sheet**:
            - Cell Membrane: Phospholipid bilayer with embedded proteins
            - Functions: Transport, signaling, protection
            - Types: Passive (diffusion) vs Active (ATP-required)
            
            **🔄 Flashcard Example**:
            - Front: "What is the primary function of mitochondria?"
            - Back: "Energy production through cellular respiration (ATP synthesis)"
            
            **📊 Study Outline**:
            - I. Cell Structure
              - A. Cell Membrane
                - 1. Composition
                - 2. Functions
            
            **🔖 Key Terms**:
            - **Osmosis**: Movement of water across semipermeable membrane
            - **Endocytosis**: Cell engulfs external material
            """)
        
        # Local AI Setup Guide
        with st.expander("🏠 Local AI Setup Guide"):
            st.markdown("""
            **Why use Local AI?**
            - ✅ **Completely Free** - No API costs ever
            - 🔐 **Private** - Your data never leaves your computer
            - 🚀 **Fast** - No internet required after setup
            - 🎯 **Always Available** - No rate limits or downtime
            
            **Quick Setup:**
            ```bash
            # 1. Install Ollama (visit ollama.ai for download)
            
            # 2. Start Ollama server
            ollama serve
            
            # 3. Download Gemma 2B model (recommended for speed)
            ollama pull gemma2:2b
            
            # 4. Or download larger models for better quality
            ollama pull gemma2:9b
            ollama pull gemma2:27b
            Or any other model you prefer
            ```
            
            **Hardware Requirements:** Examples
            - **2B Model**: 2GB RAM, runs on most computers
            - **9B Model**: 6GB RAM, better quality
            - **27B Model**: 16GB RAM, best quality
            
            **Troubleshooting:**
            - Ensure Ollama is running: `ollama list`
            - Check server status: `curl http://localhost:[PORT]/api/tags`
            - View logs: Check terminal where `ollama serve` is running
            """)
        
        # Add example of open-ended scoring
        with st.expander("🔍 See Open-ended Question Example"):
            st.markdown("""
            **Example Question:** *Explain the chemical composition and importance of water (4 marks)*
            
            **Sample Answer:** 
            "Water is a liquid encompassed with two hydrogens and one oxygen. It is a crucial composition that its existence guarantees life."
            
            **AI Scoring Breakdown:**
            - ✅ Chemical composition (H2O) - 2/2 marks
            - ⚠️ Physical properties mentioned - 1/1 mark  
            - ✅ Biological importance stated - 1/1 mark
            
            **Result:** 4/4 marks (100%) with specific feedback on scientific accuracy!
            """)
        
        st.info("💡 **Pro Tip:** The app now works gracefully even with provider errors. Configure your preferred AI provider in the sidebar and start generating quizzes or study materials!")

if __name__ == "__main__":
    main()
