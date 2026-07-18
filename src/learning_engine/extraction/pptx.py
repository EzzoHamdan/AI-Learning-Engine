"""PowerPoint text extraction (python-pptx)."""

from __future__ import annotations

import io

from pptx import Presentation


def extract_text_from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)
