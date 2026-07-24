"""Tests for the extraction package: dispatch, round-trip, and the size limit.

Documents are built in-process rather than committed as binary fixtures, so the
test data stays readable and cannot drift from the libraries that parse it.
"""

from __future__ import annotations

import io

import pytest

from learning_engine.extraction import (
    ExtractionError,
    FileTooLargeError,
    UnsupportedFormatError,
    extract_text,
)

SENTINEL = "Mitochondria are the powerhouse of the cell"


# --------------------------------------------------------------------------- #
# Document builders
# --------------------------------------------------------------------------- #


def _pdf_bytes(text: str = SENTINEL) -> bytes:
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), text)
    data = doc.tobytes()
    doc.close()
    return data


def _docx_bytes(text: str = SENTINEL) -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(text: str = SENTINEL) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


BUILDERS = {"pdf": _pdf_bytes, "docx": _docx_bytes, "pptx": _pptx_bytes}


# --------------------------------------------------------------------------- #
# Round trip
# --------------------------------------------------------------------------- #


@pytest.mark.parametrize("kind", ["pdf", "docx", "pptx"])
def test_text_survives_the_round_trip(kind):
    assert SENTINEL in extract_text(BUILDERS[kind](), kind)


@pytest.mark.parametrize("kind", ["pdf", "docx", "pptx"])
def test_extension_matching_is_case_and_dot_insensitive(kind):
    data = BUILDERS[kind]()
    for spelling in (kind.upper(), f".{kind}", f".{kind.upper()}"):
        assert SENTINEL in extract_text(data, spelling)


def test_multi_paragraph_docx_keeps_every_paragraph():
    import docx

    document = docx.Document()
    for line in ("first", "second", "third"):
        document.add_paragraph(line)
    buffer = io.BytesIO()
    document.save(buffer)

    text = extract_text(buffer.getvalue(), "docx")
    assert ["first", "second", "third"] == [ln for ln in text.splitlines() if ln.strip()]


# --------------------------------------------------------------------------- #
# Dispatch and failure modes
# --------------------------------------------------------------------------- #


@pytest.mark.parametrize("kind", ["txt", "md", "xlsx", ""])
def test_unsupported_format_is_rejected_by_name(kind):
    with pytest.raises(UnsupportedFormatError, match="Unsupported file format"):
        extract_text(b"whatever", kind)


def test_unsupported_format_error_lists_what_is_supported():
    with pytest.raises(UnsupportedFormatError, match="pdf, docx, pptx"):
        extract_text(b"whatever", "txt")


@pytest.mark.parametrize("kind", ["pdf", "docx", "pptx"])
def test_corrupt_file_raises_extraction_error_not_a_raw_parser_error(kind):
    """The UI catches ExtractionError; a leaked parser exception would crash it."""
    with pytest.raises(ExtractionError):
        extract_text(b"this is definitely not a document", kind)


# --------------------------------------------------------------------------- #
# Size limit
# --------------------------------------------------------------------------- #


def test_oversized_file_is_rejected_before_parsing():
    oversized = b"x" * (2 * 1024 * 1024)
    with pytest.raises(FileTooLargeError, match="exceeds the 1MB limit"):
        extract_text(oversized, "pdf", max_mb=1)


def test_size_limit_message_reports_the_actual_size():
    with pytest.raises(FileTooLargeError, match=r"2\.0MB"):
        extract_text(b"x" * (2 * 1024 * 1024), "pdf", max_mb=1)


def test_file_at_the_limit_is_accepted():
    data = _pdf_bytes()
    assert len(data) < 1024 * 1024
    assert SENTINEL in extract_text(data, "pdf", max_mb=1)


def test_no_limit_means_no_size_check():
    """max_mb=None is the library default; the UI supplies the configured limit."""
    with pytest.raises(ExtractionError):  # fails parsing, not the size gate
        extract_text(b"x" * (2 * 1024 * 1024), "pdf")


def test_unsupported_format_is_checked_before_size():
    """A wrong extension is the more useful error to report first."""
    with pytest.raises(UnsupportedFormatError):
        extract_text(b"x" * (2 * 1024 * 1024), "txt", max_mb=1)
